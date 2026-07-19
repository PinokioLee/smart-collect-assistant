"""최종 발표 PPT 빌드 — 템플릿의 run 텍스트를 정밀 교체(서식 보존) + 아키텍처 삽입.

수치 출처:
  data/benchmark_metrics.json         — 검출 F1 100% / 0.08초·500행+/초 / 자가교정 100%·재작업 60%↓ /
                                         재현성 표준편차 0 / 결정론 후보탐색 Strict 0.786 vs Balanced 1.0(오탐 1건 회피)
  data/llm_vs_rule_benchmark.json     — LLM 직접판단(A) vs 규칙기반(B), 46행/오류10건, LLM 5회 반복
  data/llm_vs_rule_benchmark_large.json — 동일 비교, 130행/오류26건 규모 (실무 규모 재검증)

핵심 메시지: 엑셀 검증에 LLM을 직접 써본 뒤(A), 정확도는 동일(F1 100%)하지만 처리속도가 46행 기준
약 301배(2,524ms→8.4ms)·130행 기준 약 174배(3,665ms→21ms) 느리고 API 비용·파싱 실패 위험이 있어 규칙기반 코드(B)로 전환했다.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

ROOT = Path(__file__).resolve().parent
TEMPLATE = ROOT / "template.pptx"
OUTPUT = ROOT / "AI_Master_최종발표_이형진_08079.pptx"
DIAGRAM = ROOT / "architecture.png"

MENTEE = "이형진, 08079"
MENTOR = "AI Master 멘토"
SUBJECT = "Smart Collect Assistant (엑셀 취합 자동화)"


def walk(shapes):
    """그룹 포함 모든 도형을 평탄하게 순회."""
    for sh in shapes:
        if sh.shape_type == 6:  # GROUP
            yield from walk(sh.shapes)
        else:
            yield sh


def name_map(slide):
    return {sh.name: sh for sh in walk(slide.shapes)}


def replace_runs(shape, mapping: dict[str, str]) -> None:
    """도형 내 run 텍스트를 stripped 키 기준으로 교체(서식 유지)."""
    if not shape.has_text_frame:
        return
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            key = r.text.strip()
            if key in mapping:
                r.text = mapping[key]


SPEAKER_NOTES = {
    0: (
        "안녕하세요. 취합 업무 자동화 과제 Smart Collect Assistant를 발표할 이형진입니다.\n"
        "부서별 엑셀을 사람이 일일이 취합·검증하던 일을, 멀티에이전트로 자동화한 프로젝트입니다."
    ),
    1: (
        "[프로젝트 개요 · 2분]\n"
        "먼저 문제입니다. 부서별 엑셀 취합·검증은 한 건당 약 90분이 들고, 작성기준 불일치·필수값 "
        "누락·중복·형식오류가 반복됩니다. 담당자 1인이 월 10시간 넘게 단순 반복에 쓰고, 검증 "
        "누락은 보고 신뢰도까지 떨어뜨립니다.\n"
        "그래서 취합 요청 메일을 분석해 검증기준을 만들고, 제출 엑셀을 규칙기반으로 검증→자가교정"
        "→병합하는 멀티에이전트를 만들었습니다. Requirement·Supervisor·Validation·Self-"
        "Correction·Report 다섯 에이전트가 LangGraph 위에서 협업합니다.\n"
        "성과는 세 수치입니다. ① 오류 검출 F1 100%, 오탐·미탐 0 ② 46행을 약 0.06초, 초당 500행 이상 처리 "
        "③ 자가교정으로 재작업 60% 감소.\n"
        "핵심은 이게 메일을 한 번 요약하는 도구가 아니라, LLM이 검증 전략을 계획하고 교정값을 제안하면 "
        "코드가 허용값·재검증으로 확정하는 하나의 에이전트 세션이라는 점입니다. 그 판단 과정은 요청마다 "
        "트레이스 파일로 기록되어, 주장이 아니라 기록으로 증명됩니다."
    ),
    2: (
        "[기술 아키텍처 · 4분]\n"
        "왼쪽 흐름입니다. 입력은 취합 요청 메일과 제출 엑셀. Requirement Analysis가 메일에서 "
        "작성항목·마감을 추출하고, Supervisor가 검증규칙을 정한 뒤, Excel Validation이 규칙기반 "
        "검증, Self-Correction이 자동교정, Report가 요약합니다. 전체를 LangGraph StateGraph가 "
        "상태 기반으로 오케스트레이션합니다.\n"
        "기술 선택 셋. ① LangGraph — 단일 에이전트·단순 Chain으로는 분석·검증·교정·보고의 책임 "
        "분리와 상태 관리가 어렵습니다. 조건분기로 7개 노드를 오케스트레이션했습니다.\n"
        "② Supervisor는 LLM이 실제 업로드 컬럼을 보고 검증 전략과 리스크를 계획합니다. 예를 들어 "
        "'긴급도 코드값이 흔들릴 수 있다'를 미리 지목합니다. 규칙 선택 자체는 ToT 탐색 패턴에서 착안한 후보 "
        "3개를 만들어 결정론적으로 고릅니다 — 항목을 전부 필수로 강제(Strict)하면 드리프트에서 오탐이 나서 "
        "(실측 Strict 0.786점·오탐 1건 vs Balanced 1.0점·오탐 0건) Balanced를 택했습니다. Yao 2023 참고.\n"
        "③ Self-Correction이 이 프로젝트의 에이전틱 핵심입니다. 고칠 수 있는 오류에 대해 LLM이 교정값을 "
        "근거와 함께 제안하고 — 예: 긴급도 '매우 급함'을 '상'으로, 근거는 '의미가 상과 가장 유사' — 코드가 "
        "허용값·날짜형식으로 재검증해 통과한 제안만, 그리고 재검증에서 오류가 줄 때만 채택합니다. LLM 제안이 "
        "게이트를 통과 못하면 규칙으로 폴백합니다. Madaan 2023 참고.\n"
        "④ 검증·병합 자체는 규칙기반입니다. LLM 직접 판단과 실측 비교하면 정확도는 F1 100%로 같지만 규칙기반이 "
        "174~301배 빠르고 재현성 표준편차 0입니다. 그리고 이 모든 판단은 요청마다 트레이스 파일로 기록됩니다."
    ),
    3: (
        "[핵심 기술 과제 · 4분]\n"
        "가장 중요하게 보여드리고 싶은 부분입니다. 엑셀 검증을 LLM이 직접 판단하게 할지, 규칙기반 "
        "코드로 할지를 두고 실제로 두 방식을 실측 비교했습니다.\n"
        "먼저 A안, Azure OpenAI(gpt-4.1)로 행을 직접 판단시켰습니다. 정확도는 F1 100%로 준수했지만 "
        "문제는 속도였습니다. 같은 라벨 데이터셋으로 5회씩 반복했더니 46행 기준 평균 2,524밀리초, "
        "130행 기준 평균 3,665밀리초가 걸렸습니다. 호출마다 API 비용이 발생하고, 응답이 JSON "
        "형식을 벗어나면 파싱에 실패할 위험도 있었습니다.\n"
        "그래서 B안, pandas·openpyxl 규칙기반 코드로 검증·병합을 옮겼습니다. 정확도는 똑같이 F1 "
        "100%인데 처리시간은 46행 8.4밀리초, 130행 21밀리초로 174~301배 빨랐고, 10회 반복해도 "
        "표준편차 0으로 완전히 결정론적이었습니다. 비용도 추가로 들지 않습니다.\n"
        "두 번째 사례는 검증 규칙 선택입니다. 메일에서 뽑은 항목을 전부 필수컬럼으로 강제하는 "
        "Strict 방식은 양식이 부서마다 다르면 오탐을 냈습니다 — 점수 0.786, 오탐 컬럼 1개. Tree of "
        "Strict·Balanced·Loose 코드 정의 후보 3개를 실제 파일과 비교 평가하는 결정론 방식으로 "
        "바꾸자 점수 1.0, 오탐 0건이 됐습니다.\n"
        "정리하면, 정확도가 같다면 비용·속도·재현성이 유리한 쪽을 선택해야 한다는 것이 이번 과제의 "
        "핵심 판단이었습니다. 그래서 판단이 필요한 곳 — 검증 전략 계획과 교정값 제안 — 에는 LLM을 쓰고, "
        "정확성과 재현성이 필요한 검증·확정은 규칙기반 코드에 맡겼습니다. 그리고 이 판단 과정 전체를 "
        "요청마다 트레이스 파일로 기록해, 어느 단계를 LLM이 판단했고 코드가 검증했는지 증명할 수 있게 "
        "했습니다. 감사합니다."
    ),
}


def set_notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text


def main() -> None:
    prs = Presentation(str(TEMPLATE))
    slides = prs.slides

    # ---------- 표지 (slide 1) ----------
    nm = name_map(slides[0])
    replace_runs(nm["Text 4"], {"[여기에 과제명을 입력하세요]": SUBJECT})
    replace_runs(nm["Text 5"], {"[성명], [사번]": MENTEE})
    replace_runs(nm["Text 6"], {"[성명1], [성명2]": MENTOR})

    # ---------- 슬라이드 1: 프로젝트 개요 (slide 2) ----------
    nm = name_map(slides[1])
    replace_runs(nm["Text 10"], {
        "예: 분산된 사내 지식으로 인한 업무 비효율성 해결 및 반복 보고 업무 자동화":
            "부서별 엑셀 취합·검증을 수작업에 의존 — 1건당 약 90분, 작성기준 불일치·필수값 누락·중복·형식오류가 반복 발생.",
        "해결되지 않으면 어떤 손실이 발생하는지 간략히 서술하세요.":
            "취합 담당 1인이 월 10시간+를 단순 반복에 소모하고, 검증 누락이 최종 보고 신뢰도를 저하시킵니다.",
    })
    replace_runs(nm["Text 14"], {
        "예: RAG 기반 지식 검색 시스템 + 멀티 에이전트 워크플로우":
            "LLM이 검증전략을 계획하고 교정을 제안하면, 코드가 검증·확정하는 멀티에이전트 세션 (분석→검증→자가교정→병합).",
        "Agent 1 [역할]  →  Agent 2 [역할]  →  Agent 3 [역할]":
            "Requirement → Supervisor(LLM계획) → Validation → Self-Correction(LLM제안) → Report",
        "LangGraph / MCP / Azure OpenAI / RAG Pipeline 등":
            "LangGraph · Azure OpenAI · pandas/openpyxl · FastAPI · React",
    })
    # 핵심 성과 수치 3개
    replace_runs(nm["Text 17"], {
        "[수치 1]": "F1 100%", "성과 지표 설명": "오류 검출 정확도",
        "정성적 지표": "오탐 0·미탐 0",
    })
    replace_runs(nm["Text 19"], {
        "[수치 2]": "약 0.06초", "성과 지표 설명": "46행 처리",
        "정량적 지표": "500행+/초·로드+검증+병합",
    })
    replace_runs(nm["Text 21"], {
        "[수치 3]": "60%↓",
        "성과 지표 설명 (예: 업무 자동화율, 정보 검색 시간 단축 등)":
            "자가교정으로 재작업 감소 (교정율 100%)",
    })
    replace_runs(nm["Text 24"], {
        '"프로젝트의 핵심 가치를 한 문장으로 정의하세요.':
            '"LLM은 검증전략을 계획하고 교정을 제안, 코드는 허용값·재검증으로 확정 —',
        '예: RAG와 멀티 에이전트 기술로 업무 정보 검색 시간을 80% 단축했습니다."':
            'F1 100%·재현성 표준편차 0으로 자동화하고, 그 판단 과정을 트레이스로 기록·증명합니다."',
    })

    # ---------- 슬라이드 2: 기술 아키텍처 (slide 3) ----------
    nm = name_map(slides[2])
    replace_runs(nm["Text 14"], {"기술명 (예: LangGraph + MCP)": "LangGraph Multi-Agent Supervisor"})
    replace_runs(nm["Text 15"], {
        "단순 Chain 방식으로 해결되지 않는 이유와 이 기술이 필요한 근거를 서술하세요.":
            "단일 Agent/단순 Chain은 분석·검증·교정·보고의 책임 분리와 상태 관리가 어렵습니다.",
        "어떻게 구현했는지 핵심 포인트를 1~2줄로 작성하세요.":
            "State 기반 조건분기로 7개 노드를 오케스트레이션하고, 기준 검색은 선택 기능으로 분리했습니다.",
    })
    replace_runs(nm["Text 20"], {"기술명 (예: RAG + Vector DB)": "후보 탐색 + Self-Correction (LLM 제안→코드 검증)"})
    replace_runs(nm["Text 21"], {
        "Fine-tuning, 키워드 검색 등 대안 대비 RAG를 선택한 근거를 서술하세요.":
            "LLM이 검증전략 계획·교정값을 제안하면, 코드가 허용값·재검증 게이트로 확정(환각 차단).",
        "검색 정확도, 응답 속도 등 측정 가능한 개선 지표를 포함하세요.":
            "결정론 후보선택 0.786→1.0(오탐 1→0건) · 교정은 재검증서 오류 줄 때만 채택 (Yao·Madaan 2023에서 착안).",
    })
    replace_runs(nm["Text 26"], {"기술명 (예: Azure OpenAI / Whisper)": "규칙기반 검증 (pandas/openpyxl)"})
    replace_runs(nm["Text 27"], {
        "보안, 성능, 비용 등 실제 선택 기준을 명시하세요.":
            "LLM 직접판단(A)도 실측했으나 규칙기반(B) 대비 46행 301배·130행 174배 느렸습니다.",
        "Prompt Engineering, Tool Calling 설계 등 핵심 내용을 작성하세요.":
            "정확도는 동일(F1 100%), 비용 0·재현성 0. 검증 함수만 46행 기준 2,524ms→8.4ms.",
    })
    # 아키텍처 다이어그램 삽입 + 플레이스홀더 제거
    diagram_box = name_map(slides[2])["Text 8"]
    replace_runs(diagram_box, {})  # noop
    for p in diagram_box.text_frame.paragraphs:
        for r in p.runs:
            r.text = ""
    slides[2].shapes.add_picture(
        str(DIAGRAM), Inches(0.45), Inches(1.35), width=Inches(4.7)
    )

    # ---------- 슬라이드 3: 핵심 기술 과제 (slide 4) ----------
    # 실측 근거: data/llm_vs_rule_benchmark.json (46행/오류10건), data/llm_vs_rule_benchmark_large.json
    # (130행/오류26건), data/benchmark_metrics.json 의 candidate_search_discrimination.
    nm = name_map(slides[3])
    replace_runs(nm["Text 9"], {
        "단순 구현으로 해결되지 않았던 핵심 기술 과제를 제시하세요.":
            "엑셀 검증을 LLM이 직접 판단하게 할 것인가(A), 규칙 기반 코드로 할 것인가(B)?",
        "예: 단순 ReAct Agent는 10단계 이상의 Long-horizon task에서 계획 일관성(Coherence)을 잃고 실패했습니다.":
            "① LLM 직접판단(A): F1 100%지만 규칙기반 대비 46행 301배·130행 174배 느림, API비용·파싱실패 위험  ② 항목 전부 필수컬럼 강제(Strict 방식): 드리프트 시 오탐 발생",
    })
    replace_runs(nm["Text 13"], {
        "예: Planner Agent가 먼저 전체 계획을 시뮬레이션(MCP 핵심)하고, LangGraph Self-Correction 루프로 계획을 동적 수정하는 Deep Reasoning 구조를 설계했습니다.":
            "LLM은 판단(계획·제안), 코드는 검증·확정으로 분리. ① 검증: LLM 직접판단(A) 실측 후 규칙기반(B) 전환 ② 규칙선택: 결정론 후보 비교 ③ 자가교정: LLM 제안→코드 검증 후 채택",
        "Tree of Thoughts / Self-Correction / Deep Reasoning 등 고급 기법을 사용했다면 논문 레퍼런스와 함께 적용 방법을 명시하세요.":
            "실측: LLM 3,665ms(130행)·2,524ms(46행) vs 규칙기반 21ms·8.4ms → 174~301배. 후보탐색: Strict 0.786→Balanced 1.0(오탐 1→0건). 참고: Yao 2023·Madaan 2023에서 착안",
        "단순 방법 대비 이 설계가 필요했던 기술적 근거를 1~2줄로 서술하세요.":
            "정확도가 같으면 비용·속도·재현성이 유리한 쪽을 택함 — LLM은 호출마다 과금·지연 편차·네트워크 의존이 있지만, 규칙기반은 무료·고속·표준편차 0. "
            "그래서 판단(계획·제안)은 LLM, 확정(검증·병합)은 코드로 나눴고, 이 판단 과정은 요청마다 트레이스(.md/.json)로 자동 기록해 ‘주장’이 아니라 ‘기록’으로 증명합니다.",
    })
    # 결과 및 성과 (그룹 내 stat 3개) — LLM(A) vs 규칙기반(B) 실측 비교 + 후보탐색 실측 비교
    nm = name_map(slides[3])
    replace_runs(nm["Text 18"], {"XX%": "174배+"})
    replace_runs(nm["Text 20"], {
        "성과 지표 1": "검증 단계 처리속도",
        "개선 내용 예: API 호출 성공률 향상": "LLM(A) vs 규칙기반(B) 130행 실측 · 46행은 301배",
    })
    replace_runs(nm["Text 22"], {"XX%": "1.0"})
    replace_runs(nm["Text 24"], {
        "성과 지표 2": "후보 탐색 규칙 선택 점수",
        "개선 내용 예: 태스크 자동화율 달성": "Strict 0.786→Balanced 1.0, 오탐 1건→0건",
    })
    replace_runs(nm["Text 26"], {"XX배": "0"})
    replace_runs(nm["Text 28"], {
        "성과 지표 3": "재현성 표준편차",
        "개선 내용 예: 처리 속도 또는 정확도 향상": "규칙기반 10회 반복 결정론 (LLM은 지연 편차 존재)",
    })

    # 발표자 노트 = 10분 발표 대본으로 교체
    for idx, note in SPEAKER_NOTES.items():
        set_notes(slides[idx], note)

    prs.save(str(OUTPUT))
    print("saved:", OUTPUT)


if __name__ == "__main__":
    main()
