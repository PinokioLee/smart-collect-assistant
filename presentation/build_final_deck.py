# -*- coding: utf-8 -*-
"""최종 발표 덱 빌드 (표지+3장). 개인 동기·정량 임팩트·아키텍처 다이어그램 포함.

수치 출처(현재 기준):
  data/roi_benchmark_llm.json  — Fixed 78.6% → Agentic 100%, 구조실패복구 0→100%(3/3), unsafe 0 (n=14)
  data/classifier_benchmark.json — 휴리스틱 45.8% → Azure LLM 100% 의도일치(24셋), 중앙 2.46초
  data/llm_vs_rule_benchmark_large.json — 규칙 F1 1.0/21.12ms vs Direct LLM F1 1.0/3,934ms (130행·26오류, ~186배), 재현성 stdev 0
  사람 대비: 취합 담당 1인 기준 주 5h+·연 200h+ (발표자 현업 경험 추정, 실측 아님)
"""
from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parent
OUT = ROOT / "08079_이형진_최종발표.pptx"

# ---- 팔레트 (네이비/민트/회색) ----
NAVY = RGBColor(0x1E, 0x33, 0x55)
NAVY2 = RGBColor(0x2B, 0x45, 0x6E)
MINT = RGBColor(0x2F, 0xB9, 0xB2)
MINT_BG = RGBColor(0xE7, 0xF7, 0xF6)
GRAY_BG = RGBColor(0xF2, 0xF4, 0xF7)
GRAY_LINE = RGBColor(0xD5, 0xDA, 0xE1)
INK = RGBColor(0x22, 0x2B, 0x38)
SUB = RGBColor(0x5A, 0x66, 0x76)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "맑은 고딕"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, l, t, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE, radius=None):
    sp = s.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def text(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True, space=1.0):
    """runs: list of paragraphs; each paragraph = list of (txt,size,color,bold) tuples."""
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(4)
    tf.margin_top = tf.margin_bottom = Pt(2)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = space
        p.space_after = Pt(2)
        for (txt, size, color, bold) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.name = FONT
    return tb


def title_bar(s, label, page):
    rect(s, 0.5, 0.36, 0.12, 0.5, fill=NAVY)                    # 좌측 강조바
    text(s, 0.72, 0.32, 8.5, 0.6, [[(label, 26, INK, True)]])
    # 권장시간 pill
    pill = rect(s, 10.9, 0.42, 1.9, 0.42, fill=GRAY_BG, line=GRAY_LINE, line_w=0.75,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    tf = pill.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "권장 " + page; r.font.size = Pt(11); r.font.color.rgb = SUB; r.font.name = FONT
    rect(s, 0.5, 1.02, 12.33, 0.02, fill=GRAY_LINE)            # 구분선


def page_no(s, n):
    text(s, 11.8, 7.02, 1.4, 0.35, [[(f"{n} / 03", 11, SUB, False)]], align=PP_ALIGN.RIGHT)


def chevron(s, cx, cy):
    text(s, cx - 0.3, cy, 0.6, 0.22, [[("▼", 12, MINT, True)]], align=PP_ALIGN.CENTER)


# =========================================================================
# SLIDE 1 — 표지
# =========================================================================
s = slide()
rect(s, 0, 0, Emu(SW).inches, 0.28, fill=NAVY)
rect(s, 0, Emu(SH).inches - 0.28, Emu(SW).inches, 0.28, fill=NAVY)
text(s, 0, 1.7, 13.333, 0.5, [[("AI Master Project 3기", 20, SUB, False)]], align=PP_ALIGN.CENTER)
text(s, 0, 2.35, 13.333, 1.2, [[("Smart Collect", 60, INK, True)]], align=PP_ALIGN.CENTER)
rect(s, 6.07, 3.75, 1.2, 0.07, fill=MINT)                      # 민트 언더라인
text(s, 0, 4.2, 13.333, 0.5,
     [[("취합 업무를 대신 수행하는 자율 멀티 에이전트", 18, NAVY2, True)]], align=PP_ALIGN.CENTER)
info = [
    [("과제명   ", 15, SUB, True), ("Gmail 이벤트 기반 자율 취합 멀티 에이전트", 15, INK, False)],
    [("발표자   ", 15, SUB, True), ("이형진 · 08079", 15, INK, False)],
    [("멘  토   ", 15, SUB, True), ("정영필, 김찬우", 15, INK, False)],
]
text(s, 0, 4.95, 13.333, 1.3, info, align=PP_ALIGN.CENTER, space=1.25)
text(s, 0, 6.75, 13.333, 0.4,
     [[("발표 영상 10분 이내   |   시연 영상 별도 제출 5분 이내", 12, SUB, False)]], align=PP_ALIGN.CENTER)

# =========================================================================
# SLIDE 2 — 프로젝트 개요
# =========================================================================
s = slide()
title_bar(s, "프로젝트 개요", "2분 이내")
page_no(s, "01")

# ---- 좌측: 문제 정의 ----
rect(s, 0.5, 1.2, 6.55, 2.75, fill=GRAY_BG, line=GRAY_LINE, line_w=0.75)
rect(s, 0.5, 1.2, 0.09, 2.75, fill=NAVY)
text(s, 0.72, 1.32, 6.2, 0.4, [[("1. 문제 정의 ", 15, NAVY, True), ("(Vision / Problem)", 12, SUB, True)]])
prob = [
    [("어떤 문제를 해결하는가?", 12, MINT, True)],
    [("취합 업무는 담당자의 본업이 아닌 ", 11.5, INK, False), ("‘부수 업무’", 11.5, INK, True),
     ("입니다. 매번 낯선 취합 주제를 새로 파악해 안내문을 쓰고, 부서별 회신을 받아 검증하고, "
      "미제출자를 일일이 독촉해야 합니다.", 11.5, INK, False)],
    [("왜 중요한가? (개인 동기 · 임팩트)", 12, MINT, True)],
    [("발표자는 ", 11.5, INK, False), ("팀 기획 업무", 11.5, INK, True),
     (" 담당자입니다. 본업 외 취합 부수업무에 시간을 뺏겨 야근까지 하는 현실의 어려움을, "
      "AI로 직접 해결하고자 시작했습니다.", 11.5, INK, False)],
    [("담당 1인 기준 ", 11.5, NAVY, True), ("주 5시간+ · 연 200시간+", 12.5, NAVY, True),
     ("  소모(현업 기준 추정)", 10.5, SUB, False)],
]
text(s, 0.72, 1.78, 6.2, 2.1, prob, space=1.05)

# ---- 좌측: 핵심 기능 ----
rect(s, 0.5, 4.08, 6.55, 2.95, fill=MINT_BG, line=GRAY_LINE, line_w=0.75)
rect(s, 0.5, 4.08, 0.09, 2.95, fill=MINT)
text(s, 0.72, 4.2, 6.2, 0.4, [[("2. 핵심 기능 ", 15, NAVY, True), ("(Core Features)", 12, SUB, True)]])
feat = [
    [("수신함을 ", 11.5, INK, False), ("상시 모니터링", 11.5, INK, True),
     ("해 메일이 올 때마다 자동 분류 (일반·취합·스팸)", 11.5, INK, False)],
    [("취합 요청이면 ", 11.5, INK, False), ("내용·첨부 양식을 파악", 11.5, INK, True),
     (" → 안내 메일 작성·발송", 11.5, INK, False)],
    [("회신을 분류해 ", 11.5, INK, False), ("정상/오류 판정", 11.5, INK, True),
     (" → 오류면 자동 재작성 요청 메일 발송", 11.5, INK, False)],
    [("정상 회신은 ", 11.5, INK, False), ("작성완료자 명단에 등록", 11.5, INK, True),
     (", 미제출자는 자동 리마인드", 11.5, INK, False)],
    [("실패·예외는 ", 11.5, INK, False), ("사람 승인", 11.5, INK, True),
     ("으로 안전하게 전환 (Human-in-the-loop)", 11.5, INK, False)],
    [("기술 스택  ", 10.5, SUB, True),
     ("LangGraph · Azure OpenAI · Gmail API · pandas/openpyxl · SQLite · Langfuse", 10.5, INK, False)],
]
text(s, 0.72, 4.66, 6.2, 2.3, feat, space=1.12)

# ---- 우측: 핵심 성과 3타일 ----
text(s, 7.35, 1.18, 5.4, 0.35, [[("핵심 성과 ", 15, INK, True), ("(Key Achievements)", 12, SUB, True)]])
tiles = [
    ("연 200시간+ → 자동화", "취합 반복업무를 에이전트가 대행 · 담당자는 승인만 (현업 기준)", NAVY),
    ("업무 성공률 78.6 → 100%", "동일 기능 14개 시나리오 · Fixed 대비 +21.4%p (Live LLM)", MINT),
    ("구조 실패 복구 0 → 100%", "작업번호 없음·첨부 유실·손상 파일 3/3 안전 복구", NAVY),
]
ty = 1.62
for head, sub, accent in tiles:
    rect(s, 7.35, ty, 5.45, 0.92, fill=GRAY_BG, line=GRAY_LINE, line_w=0.75)
    rect(s, 7.35, ty, 0.09, 0.92, fill=accent)
    text(s, 7.6, ty + 0.13, 5.1, 0.4, [[(head, 16, NAVY, True)]])
    text(s, 7.6, ty + 0.54, 5.1, 0.35, [[(sub, 10.5, SUB, False)]])
    ty += 1.03

# ---- 우측: Key Message ----
rect(s, 7.35, 4.78, 5.45, 2.25, fill=NAVY)
text(s, 7.6, 4.95, 5.0, 0.4, [[("Key Message", 13, MINT, True)]])
text(s, 7.6, 5.45, 5.0, 1.5,
     [[("취합 담당자가 ‘본업처럼’ 매달리던 반복 업무를,", 14, WHITE, True)],
      [("LLM이 판단하고 코드가 검증하는", 14, WHITE, True)],
      [("자율 에이전트가 대신 수행한다.", 14, WHITE, True)]], space=1.25)

# =========================================================================
# SLIDE 3 — 기술 아키텍처 (다이어그램)
# =========================================================================
s = slide()
title_bar(s, "기술 아키텍처", "4분 이내")
page_no(s, "02")

# ---- 좌측 다이어그램 프레임 ----
DX, DW = 0.5, 6.5
rect(s, DX, 1.2, DW, 5.83, fill=WHITE, line=GRAY_LINE, line_w=1.0)
# 범례
text(s, DX + 0.2, 1.3, DW - 0.4, 0.3,
     [[("🧠 LLM 판단", 10.5, MINT, True), ("      ", 10.5, INK, False),
       ("⚙️ 코드 검증·확정", 10.5, NAVY, True)]], align=PP_ALIGN.CENTER)

cx = DX + DW / 2          # 다이어그램 중앙 x
bw = 5.4                  # 박스 폭
bl = cx - bw / 2


def dbox(top, title, desc, kind):
    """kind: 'llm'(민트) / 'code'(네이비) / 'io'(회색)"""
    if kind == "llm":
        fill, tcol, dcol = MINT_BG, NAVY, SUB
        rect(s, bl, top, 0.08, 0.62, fill=MINT)
    elif kind == "code":
        fill, tcol, dcol = NAVY, WHITE, RGBColor(0xC9, 0xD6, 0xE5)
    else:
        fill, tcol, dcol = GRAY_BG, INK, SUB
    rect(s, bl, top, bw, 0.62, fill=fill, line=GRAY_LINE, line_w=0.75)
    text(s, bl + 0.22, top + 0.07, bw - 0.4, 0.28, [[(title, 12.5, tcol, True)]])
    text(s, bl + 0.22, top + 0.34, bw - 0.4, 0.24, [[(desc, 9.5, dcol, False)]])


y = 1.72
dbox(y, "Gmail 수신함 · APScheduler", "새 메일 자동 수집 (스케줄 · 지금 실행)", "io"); chevron(s, cx, y + 0.64); y += 0.92
dbox(y, "🧠 Intake Agent", "일반 · 취합 · 스팸 3분류 + 5개 업무 의도", "llm"); chevron(s, cx, y + 0.64); y += 0.92
dbox(y, "🧠 Supervisor LLM  +  ⚙️ Policy Gate", "상태를 보고 다음 Worker 선택 · 발송 안전성 검증", "llm"); chevron(s, cx, y + 0.64); y += 0.92
dbox(y, "⚙️ Worker + 🧠 Self-Correction", "양식·안내 / 검증·자가교정 / Q&A / 리마인드", "code"); chevron(s, cx, y + 0.64); y += 0.92
dbox(y, "Observation Loop", "성공 END · 일시오류 1회 재시도 · 구조오류 Human Review", "io"); chevron(s, cx, y + 0.64); y += 0.92
dbox(y, "Agent Action Log (SQLite) + Langfuse", "판단·도구 실행 기록 + LLM 프롬프트·토큰·비용 관측", "code")

# ---- 우측 기술 3선 카드 ----
cards = [
    ("01", "LangGraph Supervisor",
     "고정 체인은 Worker 실패 후 다음 행동을 못 바꿈",
     "Fixed 78.6% → Agentic 100% · 구조복구 0→100%(3/3), n=14", NAVY),
    ("02", "Azure OpenAI · Structured Output",
     "표현이 다양한 3분류·5의도를 함께 판단",
     "휴리스틱 45.8% → Azure LLM 100% 의도일치 · 중앙 2.46초 (24셋)", MINT),
    ("03", "규칙 검증 (pandas/openpyxl)",
     "셀 규칙·발송 안전성은 재현 가능한 코드가 적합",
     "130행·26오류 동일 F1 1.0 · 규칙 21.12ms vs LLM 3,934ms (~186배)", NAVY),
]
cx2, cw = 7.35, 5.45
cy = 1.2
for num, name, why, meas, accent in cards:
    rect(s, cx2, cy, cw, 1.83, fill=GRAY_BG if accent == NAVY else MINT_BG, line=GRAY_LINE, line_w=0.75)
    rect(s, cx2, cy, 0.09, 1.83, fill=accent)
    badge = rect(s, cx2 + 0.25, cy + 0.22, 0.5, 0.42, fill=accent, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.3)
    btf = badge.text_frame; bp = btf.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run(); br.text = num; br.font.size = Pt(13); br.font.bold = True; br.font.color.rgb = WHITE; br.font.name = FONT
    text(s, cx2 + 0.95, cy + 0.22, cw - 1.2, 0.4, [[(name, 14.5, NAVY, True)]])
    text(s, cx2 + 0.95, cy + 0.72, cw - 1.2, 0.4, [[("선택 이유  ", 10, SUB, True), (why, 10.5, INK, False)]], space=1.0)
    text(s, cx2 + 0.25, cy + 1.28, cw - 0.5, 0.45, [[("측정  ", 10, MINT, True), (meas, 10.5, INK, True)]], space=1.0)
    cy += 1.96

# =========================================================================
# SLIDE 4 — 핵심 기술 과제
# =========================================================================
s = slide()
title_bar(s, "핵심 기술 과제", "4분 이내")
page_no(s, "03")

# ---- Hurdle 배너 ----
rect(s, 0.5, 1.2, 12.33, 1.05, fill=NAVY)
text(s, 0.75, 1.32, 12.0, 0.35, [[("핵심 기술 난제 ", 14, MINT, True), ("(Hurdle)", 11, RGBColor(0xC9,0xD6,0xE5), True)]])
text(s, 0.75, 1.68, 12.0, 0.5,
     [[("부수업무를 사람 손에서 떼려면 ", 12, WHITE, False), ("자율성과 안전성을 동시에", 12, WHITE, True),
       (" 확보해야 한다. 고정 체인은 실패를 관찰 못해 종료하고, LLM 단독 자동발송은 "
        "잘못된 수신자·근거 없는 답변·불완전 양식 같은 업무 사고 위험이 있다.", 12, WHITE, False)]], space=1.05)

# ---- 좌측: 접근 방법 ----
rect(s, 0.5, 2.45, 6.55, 4.58, fill=GRAY_BG, line=GRAY_LINE, line_w=0.75)
rect(s, 0.5, 2.45, 0.09, 4.58, fill=NAVY)
text(s, 0.72, 2.57, 6.2, 0.35, [[("엔지니어링 접근 방법", 14, NAVY, True)]])
appr = [
    [("• Supervisor는 ", 11, INK, False), ("enum으로 제한된 행동만", 11, INK, True), (" 제안", 11, INK, False)],
    [("• Job ID로 요청–회신–수정본을 연결", 11, INK, False)],
    [("• Worker 실패를 ", 11, INK, False), ("observation으로 되돌려 재계획", 11, INK, True)],
    [("• 일시 오류는 1회 재시도, 구조 오류·재실패는 사람 승인", 11, INK, False)],
    [("• Self-Correction: ", 11, INK, False), ("LLM 제안 → 코드 재검증 통과분만 채택", 11, INK, True),
     (" (환각 차단)", 11, INK, False)],
    [("• confidence·마감·필드·첨부·도메인·grounding을 코드 Gate로 검증", 11, INK, False)],
    [("• LLM 호출은 ", 11, INK, False), ("Langfuse 트레이스", 11, INK, True),
     ("로 프롬프트·토큰·비용까지 관측", 11, INK, False)],
]
text(s, 0.72, 3.0, 6.2, 2.6, appr, space=1.18)
rect(s, 0.72, 5.75, 6.1, 0.02, fill=GRAY_LINE)
text(s, 0.72, 5.85, 6.2, 0.9,
     [[("실패 시 실제 흐름", 11, MINT, True)],
      [("Intake → Supervisor → Validation(job_not_found)", 10, SUB, False)],
      [("→ Observation → Supervisor 재계획 → Human Review", 10, SUB, False)]], space=1.05)

# ---- 우측: 결과 3타일 ----
rect(s, 7.35, 2.45, 5.45, 4.58, fill=MINT_BG, line=GRAY_LINE, line_w=0.75)
text(s, 7.6, 2.57, 5.0, 0.35, [[("결과 및 성과 ", 14, NAVY, True), ("(Live LLM · n=14)", 11, SUB, True)]])
res = [
    ("100%", 30, "업무 시나리오 성공", "Fixed Workflow 78.6% 대비 +21.4%p"),
    ("100%", 30, "구조 실패 안전 복구 (3/3)", "Fixed 0% → Agentic 100% · 배치 지연 +40%"),
    ("200h+", 28, "연간 취합 반복업무 자동화", "주 5h+ 부수업무를 승인 중심으로 (현업 기준 추정)"),
]
ry = 3.0
for big, bsz, lab, sub in res:
    rect(s, 7.6, ry, 4.95, 1.15, fill=WHITE, line=GRAY_LINE, line_w=0.75)
    text(s, 7.72, ry + 0.2, 1.9, 0.8, [[(big, bsz, NAVY, True)]], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER, wrap=False)
    text(s, 9.7, ry + 0.24, 2.75, 0.4, [[(lab, 12.5, INK, True)]])
    text(s, 9.7, ry + 0.64, 2.75, 0.45, [[(sub, 9.5, SUB, False)]], space=1.0)
    ry += 1.28
text(s, 7.6, 6.75, 5.0, 0.3,
     [[("측정 경계  ", 9.5, SUB, True),
       ("사람 ROI는 실측 전(roi_claim_available=false), 200h는 현업 추정", 9.5, SUB, False)]])

# =========================================================================
# 발표자 노트 (10분 대본) — 현재 덱 수치 기준
# =========================================================================
NOTES = {
    0: (
        "안녕하세요. 취합 업무를 대신 수행하는 자율 멀티 에이전트, Smart Collect를 발표할 이형진입니다.\n"
        "핵심 한 줄: LLM이 다음 행동을 판단하고, 코드 정책이 그 행동의 실행 가능 여부를 검증하는 에이전트입니다."
    ),
    1: (
        "[프로젝트 개요 · 2분]\n"
        "먼저 왜 이 주제를 정했는지입니다. 취합 업무는 담당자의 본업이 아닌 부수 업무입니다. 그런데도 "
        "매번 낯선 취합 주제를 새로 파악해 안내문을 쓰고, 부서별 회신을 받아 검증하고, 미제출자를 일일이 "
        "독촉해야 합니다.\n"
        "저는 팀 기획 업무를 하는데, 본업 외에 이 취합 부수업무에 시간을 뺏겨 야근까지 하게 됐습니다. "
        "담당 1인 기준 주 5시간 이상, 연 200시간 이상이 여기에 들어갑니다. 제 실제 어려움을 AI로 직접 "
        "해결하려고 시작한 프로젝트입니다.\n"
        "만든 것은 이렇습니다. 수신함을 상시 모니터링해 메일이 올 때마다 자동 분류하고, 취합 요청이면 "
        "내용과 첨부 양식을 파악해 안내 메일을 작성·발송합니다. 회신이 오면 정상인지 오류인지 판정해서, "
        "오류면 자동으로 재작성 요청을 보내고, 정상이면 작성완료자 명단에 등록하며, 미제출자는 리마인드합니다.\n"
        "성과는 세 가지입니다. 연 200시간 이상 걸리던 반복 업무를 에이전트가 대행해 담당자는 승인만 하고, "
        "동일 기능 비교에서 업무 성공률이 78.6%에서 100%로, 구조적 실패 복구가 0%에서 100%로 올라갔습니다."
    ),
    2: (
        "[기술 아키텍처 · 4분]\n"
        "왼쪽 다이어그램입니다. 민트색은 LLM이 판단하는 단계, 네이비는 코드가 검증·확정하는 단계입니다. "
        "Gmail 수신함을 수집하면 Intake Agent가 일반·취합·스팸 3분류와 5개 업무 의도를 판단하고, "
        "Supervisor LLM이 현재 상태를 보고 다음 Worker를 고릅니다. Policy Gate가 발송 안전성을 코드로 "
        "검증합니다. Worker 실패는 Observation Loop로 되돌아가 재계획되고, 모든 판단은 Agent Action Log와 "
        "Langfuse 트레이스에 남습니다.\n"
        "기술 선택 세 가지입니다. 첫째 LangGraph Supervisor — 고정 체인은 Worker 실패 후 다음 행동을 바꾸지 "
        "못합니다. 동일 기능 14개 시나리오에서 고정 워크플로 78.6% 대비 Agentic 100%였고, 구조 실패 복구는 "
        "0%에서 100%로 개선됐습니다.\n"
        "둘째 Azure OpenAI Structured Output — 표현이 다양한 3분류·5의도를 함께 판단해야 해서 썼습니다. "
        "우회 표현 24개 평가셋에서 휴리스틱 45.8% 대비 Live LLM 100% 의도일치였고, 중앙 응답시간은 2.46초입니다.\n"
        "셋째 규칙 검증 — 셀 규칙과 발송 안전성은 재현 가능한 코드가 적합합니다. 130행·26오류 검증에서 규칙과 "
        "LLM 모두 F1 1.0이지만 규칙은 21밀리초, LLM은 3,934밀리초로 약 186배 빠르고 재현 편차가 0입니다. "
        "그래서 판단은 LLM, 검증·확정은 코드로 나눴습니다."
    ),
    3: (
        "[핵심 기술 과제 · 4분]\n"
        "가장 어려운 과제는 자율성과 안전성을 동시에 확보하는 것이었습니다. 부수업무를 사람 손에서 떼려면 "
        "잘못된 수신자, 근거 없는 답변, 불완전한 양식 없이 자동 처리해야 하는데, 고정 체인은 실패를 관찰하지 "
        "못하고 그냥 종료합니다.\n"
        "그래서 이렇게 접근했습니다. Supervisor 출력은 enum으로 제한된 행동만 제안하게 하고, 요청과 회신을 "
        "Job ID로 연결했습니다. Worker 실패는 observation으로 되돌려 재계획하되, 일시 오류는 1회 재시도, 구조 "
        "오류나 재실패는 사람 승인으로 보냅니다. 자가교정은 LLM이 근거와 함께 제안하면 코드가 재검증해 통과한 "
        "제안만, 그것도 오류가 줄 때만 채택합니다. confidence·마감·필드·첨부·도메인·grounding은 코드 Gate가 "
        "검증하고, 모든 LLM 호출은 Langfuse로 프롬프트·토큰·비용까지 관측합니다.\n"
        "예를 들어 작업번호 없는 제출이 오면, Intake에서 Supervisor로, Validation에서 job_not_found로 실패하고, "
        "그게 Observation으로 되돌아가 Supervisor가 재계획해 Human Review로 안전하게 넘어갑니다. 이게 고정 "
        "파이프라인과의 결정적 차이입니다.\n"
        "결과는 업무 성공률 100%, 구조 실패 복구 100%로 3건 모두 복구, 그리고 연 200시간 이상 걸리던 취합 "
        "반복업무를 승인 중심으로 자동화했습니다. 다만 사람 대비 시간 ROI는 아직 스톱워치 실측 전이라 "
        "roi_claim_available은 false이고, 200시간은 현업 경험 기준 추정임을 분명히 밝힙니다. LLM은 판단, 코드는 "
        "검증으로 역할을 나눈 자율 에이전트입니다. 감사합니다."
    ),
}
for idx, note in NOTES.items():
    prs.slides[idx].notes_slide.notes_text_frame.text = note

prs.save(str(OUT))
print("saved:", OUT, "(발표자 노트 4개 포함)")
